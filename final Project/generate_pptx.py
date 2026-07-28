"""
NTB Event Management System — Midterm Defense Presentation
Professional dark-themed PPTX with NTB Nepal colours (deep green + crimson + gold)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1F, 0x13)   # very dark forest green  #0D1F13
C_BG2       = RGBColor(0x12, 0x2B, 0x1A)   # slightly lighter green  #122B1A
C_ACCENT    = RGBColor(0xC8, 0x10, 0x2E)   # NTB crimson             #C8102E
C_GOLD      = RGBColor(0xD4, 0xAF, 0x37)   # NTB gold                #D4AF37
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xF0, 0xEA)   # near-white green tint
C_MUTED     = RGBColor(0x9A, 0xB0, 0x9E)   # muted text
C_CARD      = RGBColor(0x1A, 0x35, 0x22)   # card background
C_DIVIDER   = RGBColor(0x2A, 0x4A, 0x33)   # subtle divider

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ════════════════════════════════════════════════════════════════════════════
# Low-level helpers
# ════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, size=18, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return txb

def add_para(tf, text, size=16, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, space_before=Pt(4), font="Calibri"):
    from pptx.util import Pt as Pt2
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt2(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return p

def set_bg(slide, color=C_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def full_bg(slide):
    """Paint the full slide background dark green."""
    set_bg(slide, C_BG)

def accent_bar(slide, y=Inches(0), h=Inches(0.06)):
    """Thin crimson accent line."""
    add_rect(slide, 0, y, SLIDE_W, h, fill=C_ACCENT)

def gold_bar(slide, y=Inches(0), h=Inches(0.06)):
    add_rect(slide, 0, y, SLIDE_W, h, fill=C_GOLD)

def slide_number(slide, num, total=18):
    txt = f"{num} / {total}"
    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.35),
                Inches(1.1), Inches(0.3), txt,
                size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)

def footer_line(slide):
    add_rect(slide, 0, SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.003), fill=C_DIVIDER)
    add_textbox(slide, Inches(0.4), SLIDE_H - Inches(0.42),
                Inches(6), Inches(0.35),
                "NTB Event Management System  •  Midterm Defense  •  2026",
                size=10, color=C_MUTED)

def section_tag(slide, label):
    """Small pill tag top-left."""
    add_rect(slide, Inches(0.4), Inches(0.22), Inches(2.6), Inches(0.32), fill=C_CARD)
    add_textbox(slide, Inches(0.44), Inches(0.23), Inches(2.6), Inches(0.3),
                label.upper(), size=9, bold=True, color=C_GOLD)

def slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.4), Inches(0.62), Inches(12.2), Inches(0.6),
                title, size=28, bold=True, color=C_WHITE, font="Calibri")
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(1.22), Inches(10), Inches(0.36),
                    subtitle, size=15, color=C_MUTED, font="Calibri")
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(1.0), Inches(0.045), fill=C_ACCENT)

def card(slide, x, y, w, h, title, lines, icon="", title_color=None):
    """Rounded-corner card with title + bullet lines."""
    add_rect(slide, x, y, w, h, fill=C_CARD, line_color=C_DIVIDER, line_w=Pt(1))
    ty = y + Inches(0.15)
    icon_prefix = icon + "  " if icon else ""
    add_textbox(slide, x + Inches(0.18), ty, w - Inches(0.3), Inches(0.38),
                icon_prefix + title, size=14, bold=True,
                color=title_color or C_GOLD, font="Calibri")
    ly = ty + Inches(0.42)
    lh = (h - Inches(0.65)) / max(len(lines), 1)
    for line in lines:
        add_textbox(slide, x + Inches(0.22), ly, w - Inches(0.35), max(lh, Inches(0.28)),
                    "  " + line, size=12, color=C_LIGHT, wrap=True, font="Calibri")
        ly += lh

def stat_box(slide, x, y, w, h, number, label, color=C_ACCENT):
    add_rect(slide, x, y, w, h, fill=C_CARD, line_color=color, line_w=Pt(2))
    add_textbox(slide, x, y + Inches(0.22), w, Inches(0.55),
                number, size=36, bold=True, color=color, align=PP_ALIGN.CENTER, font="Calibri")
    add_textbox(slide, x, y + Inches(0.8), w, Inches(0.38),
                label, size=12, color=C_MUTED, align=PP_ALIGN.CENTER, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  TITLE / COVER
# ════════════════════════════════════════════════════════════════════════════
def slide_01():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)

    # Left accent strip
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, fill=C_ACCENT)

    # Decorative large circle top-right
    from pptx.util import Emu
    cx = SLIDE_W - Inches(1.8)
    cy = -Inches(1.2)
    r  = Inches(4.0)
    shape = s.shapes.add_shape(9, cx, cy, r, r)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x18, 0x40, 0x28)
    shape.line.fill.background()

    # Small circle
    s.shapes.add_shape(9,
        SLIDE_W - Inches(3.8), SLIDE_H - Inches(2.2), Inches(2.2), Inches(2.2)
    ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = RGBColor(0x15, 0x33, 0x20)
    s.shapes[-1].line.fill.background()

    # Gold top bar
    gold_bar(s, y=0, h=Inches(0.055))

    # Organisation tag
    add_textbox(s, Inches(0.5), Inches(1.0), Inches(7), Inches(0.38),
                "POKHARA UNIVERSITY  •  CACS452  •  PROJECT III",
                size=11, bold=True, color=C_GOLD, font="Calibri")

    # Main title
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(9.5), Inches(1.1),
                "NTB EVENT MANAGEMENT SYSTEM",
                size=40, bold=True, color=C_WHITE, font="Calibri")

    # Subtitle
    add_textbox(s, Inches(0.5), Inches(2.65), Inches(9), Inches(0.5),
                "A Smart Event Management Platform for Nepal Tourism Board",
                size=18, italic=True, color=C_MUTED, font="Calibri")

    # Red divider
    add_rect(s, Inches(0.5), Inches(3.25), Inches(2.5), Inches(0.055), fill=C_ACCENT)

    # Meta info
    add_textbox(s, Inches(0.5), Inches(3.42), Inches(7), Inches(0.35),
                "Presented by:  Anjal Joshi  •  Year IV / Semester VIII",
                size=13, color=C_LIGHT, font="Calibri")
    add_textbox(s, Inches(0.5), Inches(3.82), Inches(7), Inches(0.35),
                "Midterm Defense  •  June 2026",
                size=13, color=C_MUTED, font="Calibri")

    # Bottom tag badges
    for i, (tag, col) in enumerate([
        ("ASP.NET Core 8", C_ACCENT),
        ("SvelteKit", C_GOLD),
        ("PostgreSQL", RGBColor(0x33, 0x6F, 0x99)),
        ("AI / TF-IDF / BM25", RGBColor(0x2E, 0x7D, 0x32)),
    ]):
        bx = Inches(0.5) + i * Inches(2.55)
        add_rect(s, bx, Inches(5.2), Inches(2.3), Inches(0.42), fill=C_CARD,
                 line_color=col, line_w=Pt(1.5))
        add_textbox(s, bx + Inches(0.1), Inches(5.22), Inches(2.1), Inches(0.36),
                    tag, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Bottom bar
    add_rect(s, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), fill=C_CARD)
    add_textbox(s, Inches(0.5), SLIDE_H - Inches(0.35), Inches(8), Inches(0.32),
                "Department of Computer Science  •  Pokhara University",
                size=10, color=C_MUTED)
    slide_number(s, 1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  AGENDA
# ════════════════════════════════════════════════════════════════════════════
def slide_02():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Agenda")
    slide_title(s, "What We'll Cover Today")
    slide_number(s, 2)

    items = [
        ("01", "Project Introduction",      "Background, problem statement & objectives"),
        ("02", "System Architecture",        "Clean Architecture + tech stack breakdown"),
        ("03", "Features Built So Far",      "Event CRUD, calendar, admin dashboard, public portal"),
        ("04", "AI Modules",                 "4 self-built algorithms — how they work"),
        ("05", "Live Demo Walkthrough",       "Key screens and user flows"),
        ("06", "Testing & Results",           "Unit tests, system tests, performance"),
        ("07", "Conclusion & Next Steps",    "What's done, what's coming"),
    ]

    for i, (num, title, desc) in enumerate(items):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.85) + row * Inches(1.08)
        add_rect(s, x, y, Inches(6.1), Inches(0.92), fill=C_CARD,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_rect(s, x, y, Inches(0.06), Inches(0.92), fill=C_ACCENT)
        add_textbox(s, x + Inches(0.2), y + Inches(0.08), Inches(0.52), Inches(0.4),
                    num, size=22, bold=True, color=C_ACCENT, font="Calibri")
        add_textbox(s, x + Inches(0.78), y + Inches(0.08), Inches(5.1), Inches(0.36),
                    title, size=15, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, x + Inches(0.78), y + Inches(0.48), Inches(5.1), Inches(0.3),
                    desc, size=11, color=C_MUTED, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3  —  PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
def slide_03():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Introduction")
    slide_title(s, "The Problem We're Solving",
                "Nepal's tourism events lack a smart, centralized digital home")
    slide_number(s, 3)

    problems = [
        ("No Central Platform",    "NTB had no single place to manage all event types — festivals, holidays, meetings, and general events — from one admin interface."),
        ("Passive Discovery",       "Users had to scroll through flat lists. No intelligence behind what gets surfaced — high-impact events were buried."),
        ("Inconsistent Tagging",    "Manual tagging by staff led to category mismatches, making filter-based discovery unreliable."),
        ("No Smart Search",         "Basic keyword search using ILIKE returned irrelevant results with no relevance ranking."),
        ("Zero Personalization",    "There was no way to recommend events related to what a visitor was currently viewing."),
    ]

    for i, (title, desc) in enumerate(problems):
        y = Inches(1.88) + i * Inches(1.02)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.88), fill=C_CARD,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_rect(s, Inches(0.4), y, Inches(0.06), Inches(0.88),
                 fill=C_ACCENT if i % 2 == 0 else C_GOLD)
        add_textbox(s, Inches(0.65), y + Inches(0.06), Inches(2.6), Inches(0.36),
                    title, size=13, bold=True, color=C_GOLD, font="Calibri")
        add_textbox(s, Inches(3.35), y + Inches(0.06), Inches(9.4), Inches(0.65),
                    desc, size=12, color=C_LIGHT, wrap=True, font="Calibri")
        # vertical divider
        add_rect(s, Inches(3.25), y + Inches(0.12), Inches(0.003), Inches(0.64),
                 fill=C_DIVIDER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4  —  OBJECTIVES
# ════════════════════════════════════════════════════════════════════════════
def slide_04():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Introduction")
    slide_title(s, "Project Objectives")
    slide_number(s, 4)

    objs = [
        ("Build",      "A production-grade event management platform for NTB with full CRUD and lifecycle management"),
        ("Implement",  "Role-based access control — Super Admin, Admin, and Client roles"),
        ("Deliver",    "An interactive dual-calendar (AD + BS Bikram Sambat) with month and list views"),
        ("Integrate",  "4 self-built AI modules: recommendations, smart search, tag suggestion, popularity ranking"),
        ("Ensure",     "Clean Architecture backend (ASP.NET Core 8) — maintainable, testable, scalable"),
        ("Produce",    "A responsive public portal for event discovery, filtering, and exploration"),
    ]

    for i, (verb, text) in enumerate(objs):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.9) + row * Inches(1.55)
        add_rect(s, x, y, Inches(6.1), Inches(1.35), fill=C_CARD,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_rect(s, x, y, Inches(6.1), Inches(0.05), fill=C_ACCENT)
        add_textbox(s, x + Inches(0.18), y + Inches(0.14), Inches(1.1), Inches(0.4),
                    verb, size=18, bold=True, color=C_ACCENT, font="Calibri")
        add_textbox(s, x + Inches(0.18), y + Inches(0.58), Inches(5.7), Inches(0.65),
                    text, size=12, color=C_LIGHT, wrap=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5  —  SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
def slide_05():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Architecture")
    slide_title(s, "System Architecture",
                "Clean Architecture backend + SvelteKit frontend")
    slide_number(s, 5)

    # Left column: layers
    layers = [
        (C_ACCENT,                  "API Layer",            "Controllers · JWT Auth · Swagger · CORS"),
        (RGBColor(0x1E, 0x5C, 0x3A), "Application Layer",  "Services · Repository Contracts · DTOs"),
        (RGBColor(0x16, 0x4A, 0x2D), "Infrastructure Layer","EF Core · Dapper · PostgreSQL · BCrypt"),
        (RGBColor(0x10, 0x3A, 0x22), "Domain Layer",        "Entities · Enums · Business Rules"),
    ]
    for i, (col, lbl, desc) in enumerate(layers):
        y = Inches(1.9) + i * Inches(1.18)
        add_rect(s, Inches(0.4), y, Inches(5.8), Inches(1.0), fill=col,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_textbox(s, Inches(0.6), y + Inches(0.08), Inches(5.2), Inches(0.38),
                    lbl, size=14, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, Inches(0.6), y + Inches(0.5), Inches(5.2), Inches(0.36),
                    desc, size=11, color=C_MUTED, font="Calibri")

    # Arrow between layers
    for i in range(3):
        ay = Inches(2.9) + i * Inches(1.18) - Inches(0.02)
        add_textbox(s, Inches(2.8), ay, Inches(0.6), Inches(0.25),
                    "▼", size=11, color=C_DIVIDER, align=PP_ALIGN.CENTER)

    # Right column: frontend + DB
    add_rect(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(1.3),
             fill=RGBColor(0x1A, 0x35, 0x25), line_color=C_GOLD, line_w=Pt(1.5))
    add_textbox(s, Inches(7.2), Inches(1.98), Inches(5.2), Inches(0.4),
                "Frontend — SvelteKit + TypeScript", size=14, bold=True, color=C_GOLD, font="Calibri")
    add_textbox(s, Inches(7.2), Inches(2.42), Inches(5.2), Inches(0.62),
                "Reactive UI  •  Admin Dashboard  •  Public Portal\nCalendar View  •  Server-Side Rendering",
                size=11, color=C_LIGHT, wrap=True, font="Calibri")

    add_rect(s, Inches(7.0), Inches(3.38), Inches(5.8), Inches(1.3),
             fill=RGBColor(0x0F, 0x2B, 0x1E), line_color=RGBColor(0x33, 0x6F, 0x99), line_w=Pt(1.5))
    add_textbox(s, Inches(7.2), Inches(3.46), Inches(5.2), Inches(0.4),
                "Database — PostgreSQL 16", size=14, bold=True,
                color=RGBColor(0x5B, 0xA4, 0xCF), font="Calibri")
    add_textbox(s, Inches(7.2), Inches(3.9), Inches(5.2), Inches(0.62),
                "EF Core for CRUD  •  Dapper for performance\nJSON columns for images & highlights",
                size=11, color=C_LIGHT, wrap=True, font="Calibri")

    add_rect(s, Inches(7.0), Inches(4.86), Inches(5.8), Inches(1.3),
             fill=RGBColor(0x2B, 0x1A, 0x1A), line_color=C_ACCENT, line_w=Pt(1.5))
    add_textbox(s, Inches(7.2), Inches(4.94), Inches(5.2), Inches(0.4),
                "AI Modules — Pure C# (No 3rd Party AI)", size=14, bold=True, color=C_ACCENT, font="Calibri")
    add_textbox(s, Inches(7.2), Inches(5.38), Inches(5.2), Inches(0.62),
                "TF-IDF  •  Cosine Similarity  •  BM25\nKeyword Extraction  •  Weighted Scoring",
                size=11, color=C_LIGHT, wrap=True, font="Calibri")

    # Arrow between frontend and backend
    add_textbox(s, Inches(6.2), Inches(3.0), Inches(0.8), Inches(0.5),
                "◄─►", size=13, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(5.9), Inches(3.18), Inches(1.3), Inches(0.3),
                "REST API", size=9, color=C_MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6  —  TECH STACK
# ════════════════════════════════════════════════════════════════════════════
def slide_06():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Architecture")
    slide_title(s, "Technology Stack")
    slide_number(s, 6)

    techs = [
        (C_ACCENT,                        "ASP.NET Core 8",      "REST API Framework",        "Controllers, middleware, JWT auth, Swagger"),
        (RGBColor(0x2E, 0x7D, 0x32),      "EF Core + Dapper",    "Data Access",               "EF Core for CRUD, Dapper for fast list queries"),
        (RGBColor(0x1A, 0x73, 0xA8),      "PostgreSQL 16",       "Database",                  "Relational DB, JSON columns, migrations"),
        (C_GOLD,                           "SvelteKit",           "Frontend Framework",        "SSR, TypeScript, reactive stores"),
        (RGBColor(0x7B, 0x1F, 0xA2),      "JWT + BCrypt",        "Security",                  "Token auth, refresh rotation, password hashing"),
        (RGBColor(0xE6, 0x5C, 0x00),      "Clean Architecture",  "Design Pattern",            "Domain → Application → Infrastructure → API"),
        (RGBColor(0x00, 0x7A, 0x6A),      "Custom AI (C#)",      "AI Modules",                "TF-IDF, BM25, Cosine Similarity — no 3rd party"),
        (RGBColor(0x5D, 0x4E, 0x37),      "Docker Ready",        "Deployment",                "Containerized backend + frontend builds"),
    ]

    for i, (col, name, category, detail) in enumerate(techs):
        c = i % 4
        r = i // 4
        x = Inches(0.35) + c * Inches(3.24)
        y = Inches(1.88) + r * Inches(2.4)
        add_rect(s, x, y, Inches(3.0), Inches(2.1), fill=C_CARD,
                 line_color=col, line_w=Pt(2))
        add_rect(s, x, y, Inches(3.0), Inches(0.08), fill=col)
        add_textbox(s, x + Inches(0.15), y + Inches(0.18), Inches(2.7), Inches(0.4),
                    name, size=15, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, x + Inches(0.15), y + Inches(0.62), Inches(2.7), Inches(0.28),
                    category, size=10, bold=True, color=col, font="Calibri")
        add_textbox(s, x + Inches(0.15), y + Inches(0.94), Inches(2.7), Inches(0.85),
                    detail, size=11, color=C_LIGHT, wrap=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7  —  FEATURES BUILT — PROGRESS
# ════════════════════════════════════════════════════════════════════════════
def slide_07():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Progress")
    slide_title(s, "What's Built — Midterm Progress",
                "Core platform complete · AI modules in active development")
    slide_number(s, 7)

    done = [
        "JWT Authentication + Refresh Token Rotation",
        "Role-based Access (SuperAdmin / Admin / Client)",
        "Full Event CRUD with Dual-Date (AD + BS)",
        "Image Upload & Management",
        "Category & Tag Management (Many-to-Many)",
        "Audit trail for Tag Associations",
        "Admin Dashboard with Pagination & Filters",
        "Public Homepage with Featured Events",
        "Interactive Calendar (Month + List View)",
        "Event Filtering (Type, Region, Category, Date)",
        "Event Lifecycle: Draft → Published → Archived",
        "Responsive UI across all screen sizes",
    ]

    inprog = [
        "Recommendation Engine (TF-IDF + Cosine) — 70%",
        "BM25 Smart Search — 60%",
        "Tag Suggestion (Keyword Extraction) — 50%",
        "Popularity Scoring Model — 80%",
    ]

    # DONE column
    add_rect(s, Inches(0.4), Inches(1.82), Inches(7.5), Inches(0.38),
             fill=RGBColor(0x1B, 0x5E, 0x20))
    add_textbox(s, Inches(0.55), Inches(1.84), Inches(7.0), Inches(0.32),
                "✓  COMPLETED", size=13, bold=True, color=C_WHITE, font="Calibri")

    for i, item in enumerate(done):
        y = Inches(2.26) + i * Inches(0.39)
        bg = C_CARD if i % 2 == 0 else RGBColor(0x16, 0x30, 0x1E)
        add_rect(s, Inches(0.4), y, Inches(7.5), Inches(0.36), fill=bg)
        add_textbox(s, Inches(0.6), y + Inches(0.04), Inches(7.1), Inches(0.3),
                    "  ●  " + item, size=11, color=C_LIGHT, font="Calibri")

    # IN PROGRESS column
    add_rect(s, Inches(8.2), Inches(1.82), Inches(4.7), Inches(0.38),
             fill=RGBColor(0x6A, 0x1A, 0x1A))
    add_textbox(s, Inches(8.35), Inches(1.84), Inches(4.3), Inches(0.32),
                "⚙  IN PROGRESS — AI MODULES", size=13, bold=True, color=C_WHITE, font="Calibri")

    for i, item in enumerate(inprog):
        y = Inches(2.26) + i * Inches(0.7)
        add_rect(s, Inches(8.2), y, Inches(4.7), Inches(0.6), fill=C_CARD,
                 line_color=C_ACCENT, line_w=Pt(1))
        # parse progress
        pct = int(item.split("—")[-1].strip().replace("%", "")) if "%" in item else 50
        name = item.split("—")[0].strip()
        add_textbox(s, Inches(8.38), y + Inches(0.04), Inches(4.1), Inches(0.28),
                    name, size=11, bold=True, color=C_GOLD, font="Calibri")
        # progress bar bg
        add_rect(s, Inches(8.38), y + Inches(0.36), Inches(4.1), Inches(0.14), fill=C_BG2)
        # progress bar fill
        bar_w = Inches(4.1 * pct / 100)
        add_rect(s, Inches(8.38), y + Inches(0.36), bar_w, Inches(0.14), fill=C_ACCENT)
        add_textbox(s, Inches(12.1), y + Inches(0.3), Inches(0.65), Inches(0.2),
                    f"{pct}%", size=9, color=C_MUTED, align=PP_ALIGN.RIGHT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8  —  EVENT DATA MODEL
# ════════════════════════════════════════════════════════════════════════════
def slide_08():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Design")
    slide_title(s, "Event Data Model",
                "Rich entity design supporting real NTB use cases")
    slide_number(s, 8)

    # Central entity box
    cx, cy, cw, ch = Inches(4.4), Inches(1.95), Inches(4.5), Inches(4.6)
    add_rect(s, cx, cy, cw, ch, fill=RGBColor(0x1C, 0x44, 0x2B),
             line_color=C_GOLD, line_w=Pt(2))
    add_rect(s, cx, cy, cw, Inches(0.45), fill=C_GOLD)
    add_textbox(s, cx + Inches(0.1), cy + Inches(0.04), cw - Inches(0.2), Inches(0.38),
                "Event  (Core Entity)", size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    fields = [
        "Title, Slug, Summary",
        "LongDescription",
        "EventType  (Festival/Meeting/Holiday/Event)",
        "Status  (Draft / Published / Archived)",
        "DateAd + EndDateAd  (AD)",
        "DateBs + EndDateBs  (BS / Bikram Sambat)",
        "Location, Region, Address",
        "Price, Rating, EntryType",
        "Images[ ]  (JSON array)",
        "Highlights[ ]  (JSON objects)",
        "Organizer + OrganizerVerified",
        "Featured, ReadTime",
    ]
    for i, f in enumerate(fields):
        fy = cy + Inches(0.55) + i * Inches(0.33)
        add_textbox(s, cx + Inches(0.15), fy, cw - Inches(0.25), Inches(0.3),
                    "  " + f, size=10, color=C_LIGHT, font="Calibri")

    # Related entities left
    rels_left = [
        ("Users", "CreatedBy / UpdatedBy\n(FK to User entity)"),
        ("Categories", "Category name\nassociated per event"),
    ]
    for i, (name, desc) in enumerate(rels_left):
        ry = Inches(2.3) + i * Inches(2.0)
        add_rect(s, Inches(0.3), ry, Inches(3.8), Inches(1.5), fill=C_CARD,
                 line_color=C_ACCENT, line_w=Pt(1.5))
        add_textbox(s, Inches(0.5), ry + Inches(0.1), Inches(3.4), Inches(0.38),
                    name, size=13, bold=True, color=C_ACCENT, font="Calibri")
        add_textbox(s, Inches(0.5), ry + Inches(0.52), Inches(3.4), Inches(0.7),
                    desc, size=10, color=C_LIGHT, wrap=True, font="Calibri")
        # connector
        add_textbox(s, Inches(4.05), ry + Inches(0.55), Inches(0.5), Inches(0.3),
                    "─►", size=12, color=C_DIVIDER)

    # Related entities right
    rels_right = [
        ("Tags  (Many-to-Many)", "Via EventTags junction\ntable with audit trail"),
        ("EventHighlights", "Stored as JSON array\nwithin event row"),
    ]
    for i, (name, desc) in enumerate(rels_right):
        ry = Inches(2.3) + i * Inches(2.0)
        add_rect(s, Inches(9.2), ry, Inches(3.8), Inches(1.5), fill=C_CARD,
                 line_color=C_GOLD, line_w=Pt(1.5))
        add_textbox(s, Inches(9.4), ry + Inches(0.1), Inches(3.4), Inches(0.38),
                    name, size=13, bold=True, color=C_GOLD, font="Calibri")
        add_textbox(s, Inches(9.4), ry + Inches(0.52), Inches(3.4), Inches(0.7),
                    desc, size=10, color=C_LIGHT, wrap=True, font="Calibri")
        add_textbox(s, Inches(8.75), ry + Inches(0.55), Inches(0.5), Inches(0.3),
                    "◄─", size=12, color=C_DIVIDER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9  —  AI OVERVIEW (ALL 4 MODULES)
# ════════════════════════════════════════════════════════════════════════════
def slide_09():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "AI Modules")
    slide_title(s, "4 Self-Built AI Modules",
                "All implemented in C# — no external AI API dependency")
    slide_number(s, 9)

    # Disclaimer badge
    add_rect(s, Inches(8.0), Inches(0.58), Inches(4.9), Inches(0.42),
             fill=RGBColor(0x6A, 0x1A, 0x1A))
    add_textbox(s, Inches(8.1), Inches(0.60), Inches(4.7), Inches(0.36),
                "Not SQL queries. Not 3rd-party AI. Pure algorithms.",
                size=11, bold=True, color=C_WHITE, font="Calibri")

    modules = [
        (C_ACCENT,
         "01  Recommendation Engine",
         "Algorithm: TF-IDF + Cosine Similarity",
         "WHERE it runs: When a user views an event detail page\nWHAT it does: Finds the 5 most content-similar events\nHOW: Builds TF-IDF vectors from title+desc+tags, computes dot-product cosine scores"),

        (C_GOLD,
         "02  Smart Search",
         "Algorithm: BM25 (Okapi Best Match 25)",
         "WHERE it runs: Public search bar\nWHAT it does: Ranks all events by relevance to the query\nHOW: Inverted index + probabilistic term weighting by frequency and doc length"),

        (RGBColor(0x2E, 0x7D, 0x32),
         "03  Auto Tag Suggestion",
         "Algorithm: TF-IDF Keyword Extraction + Levenshtein",
         "WHERE it runs: Admin event create/edit form\nWHAT it does: Suggests relevant tags as admin types title/desc\nHOW: Extracts high-weight keywords, fuzzy-matches against tag dictionary"),

        (RGBColor(0x1A, 0x73, 0xA8),
         "04  Popularity Scoring",
         "Algorithm: Weighted Multi-Factor Scoring + Exponential Decay",
         "WHERE it runs: Homepage & featured listings\nWHAT it does: Ranks events by real quality, not just date\nHOW: Score = 0.35×Rating + 0.25×Attendance + 0.20×Featured + 0.20×Recency"),
    ]

    for i, (col, title, algo, detail) in enumerate(modules):
        c = i % 2
        r = i // 2
        x = Inches(0.35) + c * Inches(6.5)
        y = Inches(1.85) + r * Inches(2.6)
        add_rect(s, x, y, Inches(6.2), Inches(2.4), fill=C_CARD,
                 line_color=col, line_w=Pt(2))
        add_rect(s, x, y, Inches(6.2), Inches(0.06), fill=col)
        add_textbox(s, x + Inches(0.18), y + Inches(0.14), Inches(5.8), Inches(0.4),
                    title, size=14, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, x + Inches(0.18), y + Inches(0.58), Inches(5.8), Inches(0.3),
                    algo, size=11, bold=True, color=col, font="Calibri")
        add_rect(s, x + Inches(0.18), y + Inches(0.92), Inches(5.7), Inches(0.002),
                 fill=C_DIVIDER)
        add_textbox(s, x + Inches(0.18), y + Inches(0.98), Inches(5.75), Inches(1.25),
                    detail, size=10.5, color=C_LIGHT, wrap=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10  —  AI DEEP DIVE: TF-IDF + COSINE
# ════════════════════════════════════════════════════════════════════════════
def slide_10():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "AI Deep Dive")
    slide_title(s, "Recommendation Engine",
                "TF-IDF vectorization → Cosine Similarity → Top-5 similar events")
    slide_number(s, 10)

    # Step flow boxes
    steps = [
        ("1", "Build Corpus",         "Concatenate title (3×) + summary (2×) + desc + category + region + tags for every event"),
        ("2", "Tokenize & IDF",        "Remove stop-words, compute IDF: log((1+N)/(1+df(t)))+1  across all event documents"),
        ("3", "TF-IDF Vectors",        "TF(t,d) × IDF(t) gives each event a sparse high-dimensional vector of term weights"),
        ("4", "Cosine Similarity",     "sim(A,B) = (A·B)/(‖A‖×‖B‖)  →  score 0–1. Sort descending, return top 5"),
    ]

    for i, (num, title, detail) in enumerate(steps):
        x = Inches(0.35) + i * Inches(3.24)
        y = Inches(1.85)
        add_rect(s, x, y, Inches(3.05), Inches(1.35), fill=C_CARD,
                 line_color=C_ACCENT, line_w=Pt(2))
        add_rect(s, x, y, Inches(3.05), Inches(0.06), fill=C_ACCENT)
        add_textbox(s, x + Inches(0.15), y + Inches(0.12), Inches(0.5), Inches(0.45),
                    num, size=26, bold=True, color=C_ACCENT, font="Calibri")
        add_textbox(s, x + Inches(0.7), y + Inches(0.15), Inches(2.2), Inches(0.38),
                    title, size=13, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, x + Inches(0.15), y + Inches(0.65), Inches(2.8), Inches(0.58),
                    detail, size=10, color=C_LIGHT, wrap=True, font="Calibri")
        if i < 3:
            add_textbox(s, x + Inches(3.05), y + Inches(0.55), Inches(0.2), Inches(0.3),
                        "►", size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Formula box
    add_rect(s, Inches(0.35), Inches(3.38), Inches(12.6), Inches(0.72),
             fill=RGBColor(0x1C, 0x1C, 0x1C), line_color=C_ACCENT, line_w=Pt(1))
    add_textbox(s, Inches(0.55), Inches(3.44), Inches(12.2), Inches(0.58),
                "TF-IDF(t,d) = TF(t,d) × IDF(t)     where     IDF(t) = log( (1+N) / (1+df(t)) ) + 1     →     sim(A,B) = (A·B) / (‖A‖ × ‖B‖)",
                size=13, color=C_GOLD, align=PP_ALIGN.CENTER, font="Calibri")

    # Why not SQL
    add_rect(s, Inches(0.35), Inches(4.28), Inches(12.6), Inches(1.75),
             fill=C_CARD, line_color=C_GOLD, line_w=Pt(1.5))
    add_textbox(s, Inches(0.55), Inches(4.36), Inches(12.0), Inches(0.38),
                "Why this is NOT a SQL query — and why that matters:",
                size=13, bold=True, color=C_GOLD, font="Calibri")

    points = [
        "All computation happens in C# memory after the event data is loaded. The database is just a store.",
        "SQL cannot compute vector dot products or cosine distances natively — this requires linear algebra in the application layer.",
        "Result: truly similar events are surfaced based on semantic content — not just shared categories or matching keywords in WHERE clauses.",
    ]
    for i, pt in enumerate(points):
        add_textbox(s, Inches(0.55), Inches(4.82) + i * Inches(0.36), Inches(12.0), Inches(0.32),
                    "  ▸  " + pt, size=11, color=C_LIGHT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11  —  AI DEEP DIVE: BM25
# ════════════════════════════════════════════════════════════════════════════
def slide_11():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "AI Deep Dive")
    slide_title(s, "Smart Search  —  BM25 Ranking",
                "Probabilistic relevance ranking, the same algorithm behind Elasticsearch")
    slide_number(s, 11)

    # Formula
    add_rect(s, Inches(0.35), Inches(1.82), Inches(12.6), Inches(0.95),
             fill=RGBColor(0x1C, 0x1C, 0x1C), line_color=C_GOLD, line_w=Pt(1.5))
    add_textbox(s, Inches(0.5), Inches(1.9), Inches(12.2), Inches(0.78),
                "Score(D, Q) = Σᵢ  IDF(qᵢ) ×  [ f(qᵢ,D) × (k₁+1) ]  /  [ f(qᵢ,D) + k₁ × (1 − b + b × |D|/avgdl) ]",
                size=14, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, font="Calibri")

    # Parameter cards
    params = [
        ("f(q,D)",   "Term Frequency",      "How many times the query term appears in the event document"),
        ("|D|",       "Document Length",     "Total word count of the event's searchable text"),
        ("avgdl",     "Avg Document Length", "Mean word count across all events — normalises doc length"),
        ("k₁ = 1.5",  "Saturation",         "Controls how much extra repeats of a term still matter"),
        ("b = 0.75",  "Length Penalty",      "Prevents long events from unfairly dominating results"),
        ("IDF(q)",    "Inverse Doc Freq",    "Rare query terms score higher — more discriminative"),
    ]
    for i, (sym, name, desc) in enumerate(params):
        c = i % 3
        r = i // 3
        x = Inches(0.35) + c * Inches(4.3)
        y = Inches(3.0) + r * Inches(1.55)
        add_rect(s, x, y, Inches(4.0), Inches(1.35), fill=C_CARD,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_textbox(s, x + Inches(0.15), y + Inches(0.1), Inches(1.1), Inches(0.52),
                    sym, size=18, bold=True, color=C_ACCENT, font="Calibri")
        add_textbox(s, x + Inches(1.35), y + Inches(0.12), Inches(2.5), Inches(0.36),
                    name, size=12, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, x + Inches(0.15), y + Inches(0.68), Inches(3.7), Inches(0.52),
                    desc, size=10.5, color=C_MUTED, wrap=True, font="Calibri")

    # vs ILIKE
    add_rect(s, Inches(0.35), Inches(6.22), Inches(12.6), Inches(0.82),
             fill=RGBColor(0x6A, 0x1A, 0x1A))
    add_textbox(s, Inches(0.55), Inches(6.28), Inches(6.0), Inches(0.35),
                "ILIKE '%dashain%'  →  returns ANY match, ranked by nothing",
                size=11, color=RGBColor(0xFF, 0x88, 0x88), font="Calibri")
    add_textbox(s, Inches(6.75), Inches(6.28), Inches(0.5), Inches(0.35),
                "vs", size=13, bold=True, color=C_MUTED, font="Calibri")
    add_textbox(s, Inches(7.3), Inches(6.28), Inches(5.4), Inches(0.35),
                "BM25  →  returns matches RANKED by actual relevance to the query",
                size=11, color=RGBColor(0x88, 0xFF, 0xAA), font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12  —  AI: TAG SUGGESTION + POPULARITY SCORING
# ════════════════════════════════════════════════════════════════════════════
def slide_12():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "AI Deep Dive")
    slide_title(s, "Tag Suggestion  &  Popularity Scoring",
                "Two more algorithms that make the platform smarter")
    slide_number(s, 12)

    # LEFT — Tag Suggestion
    add_rect(s, Inches(0.35), Inches(1.82), Inches(6.2), Inches(5.18),
             fill=C_CARD, line_color=RGBColor(0x2E, 0x7D, 0x32), line_w=Pt(2))
    add_rect(s, Inches(0.35), Inches(1.82), Inches(6.2), Inches(0.06),
             fill=RGBColor(0x2E, 0x7D, 0x32))
    add_textbox(s, Inches(0.55), Inches(1.95), Inches(5.8), Inches(0.4),
                "03  Automated Tag Suggestion", size=15, bold=True, color=C_WHITE, font="Calibri")
    add_textbox(s, Inches(0.55), Inches(2.38), Inches(5.8), Inches(0.3),
                "TF-IDF Keyword Extraction + Levenshtein Fuzzy Match",
                size=11, bold=True, color=RGBColor(0x66, 0xBB, 0x6A), font="Calibri")

    steps_tag = [
        ("1.", "Admin types event title + description in the form"),
        ("2.", "System tokenises input, removes stop words"),
        ("3.", "TF-IDF scores computed against the existing event corpus"),
        ("4.", "Top-K keywords extracted by highest TF-IDF weight"),
        ("5.", "Keywords fuzzy-matched against tag dictionary using Levenshtein distance (edit distance ≤ 2)"),
        ("6.", "Matched tags returned ranked by similarity — displayed as clickable chips in the UI"),
    ]
    for i, (num, text) in enumerate(steps_tag):
        y = Inches(2.82) + i * Inches(0.62)
        add_textbox(s, Inches(0.6), y, Inches(0.38), Inches(0.52),
                    num, size=13, bold=True, color=RGBColor(0x66, 0xBB, 0x6A), font="Calibri")
        add_textbox(s, Inches(1.05), y, Inches(5.2), Inches(0.52),
                    text, size=11, color=C_LIGHT, wrap=True, font="Calibri")

    # RIGHT — Popularity Scoring
    add_rect(s, Inches(6.85), Inches(1.82), Inches(6.1), Inches(5.18),
             fill=C_CARD, line_color=RGBColor(0x1A, 0x73, 0xA8), line_w=Pt(2))
    add_rect(s, Inches(6.85), Inches(1.82), Inches(6.1), Inches(0.06),
             fill=RGBColor(0x1A, 0x73, 0xA8))
    add_textbox(s, Inches(7.05), Inches(1.95), Inches(5.7), Inches(0.4),
                "04  Popularity Scoring Model", size=15, bold=True, color=C_WHITE, font="Calibri")
    add_textbox(s, Inches(7.05), Inches(2.38), Inches(5.7), Inches(0.3),
                "Weighted Multi-Factor Scoring + Exponential Recency Decay",
                size=11, bold=True, color=RGBColor(0x5B, 0xA4, 0xCF), font="Calibri")

    # Formula
    add_rect(s, Inches(7.05), Inches(2.82), Inches(5.7), Inches(0.6),
             fill=RGBColor(0x1C, 0x1C, 0x2C))
    add_textbox(s, Inches(7.15), Inches(2.9), Inches(5.5), Inches(0.44),
                "Score = 0.35×Rating + 0.25×Attendance + 0.20×Featured + 0.20×Recency",
                size=10.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, font="Calibri")

    weights = [
        ("0.35", "Rating Weight",      "User ratings are the strongest quality signal"),
        ("0.25", "Attendance Weight",  "Engagement/attendance level"),
        ("0.20", "Featured Boost",     "Admin-promoted events get a lift"),
        ("0.20", "Recency Decay",      "Score = e^(−0.01 × days_since_start) — fresh events rank higher"),
    ]
    for i, (w, label, desc) in enumerate(weights):
        y = Inches(3.58) + i * Inches(0.84)
        add_rect(s, Inches(7.05), y, Inches(0.72), Inches(0.62), fill=C_BG2)
        add_textbox(s, Inches(7.05), y + Inches(0.08), Inches(0.72), Inches(0.46),
                    w, size=16, bold=True, color=RGBColor(0x5B, 0xA4, 0xCF),
                    align=PP_ALIGN.CENTER, font="Calibri")
        add_textbox(s, Inches(7.82), y + Inches(0.02), Inches(4.8), Inches(0.3),
                    label, size=12, bold=True, color=C_WHITE, font="Calibri")
        add_textbox(s, Inches(7.82), y + Inches(0.34), Inches(4.8), Inches(0.28),
                    desc, size=10, color=C_MUTED, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13  —  CALENDAR & PUBLIC PORTAL SCREENS
# ════════════════════════════════════════════════════════════════════════════
def slide_13():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Demo")
    slide_title(s, "Key Features — Demo Walkthrough",
                "What the user actually sees and does")
    slide_number(s, 13)

    screens = [
        (C_ACCENT,
         "Public Homepage",
         [
             "Hero section with Nepal Tourism branding",
             "Featured events with rating + highlights",
             "Stats bar: total events, categories, regions",
             "Category-based quick navigation",
         ]),
        (C_GOLD,
         "Interactive Calendar",
         [
             "Month view with color-coded event types",
             "Click any day to see events in that day",
             "Toggle between AD (English) and BS (Nepali)",
             "Switch between Month view and List view",
         ]),
        (RGBColor(0x2E, 0x7D, 0x32),
         "Admin Dashboard",
         [
             "Event listing with pagination and filters",
             "One-click status change (Draft/Published/Archived)",
             "Category and tag management panel",
             "User management with role assignment",
         ]),
        (RGBColor(0x1A, 0x73, 0xA8),
         "Event Create / Edit Form",
         [
             "Rich fields: title, desc, type, dates, location",
             "Dual-date picker: AD and BS simultaneously",
             "Drag-and-drop image upload",
             "AI tag suggestion chips as admin types",
         ]),
    ]

    for i, (col, title, items) in enumerate(screens):
        c = i % 2
        r = i // 2
        x = Inches(0.35) + c * Inches(6.5)
        y = Inches(1.85) + r * Inches(2.55)
        add_rect(s, x, y, Inches(6.2), Inches(2.35), fill=C_CARD,
                 line_color=col, line_w=Pt(2))
        add_rect(s, x, y, Inches(6.2), Inches(0.06), fill=col)
        add_textbox(s, x + Inches(0.18), y + Inches(0.14), Inches(5.8), Inches(0.38),
                    title, size=14, bold=True, color=C_WHITE, font="Calibri")
        for j, item in enumerate(items):
            add_textbox(s, x + Inches(0.22), y + Inches(0.62) + j * Inches(0.4),
                        Inches(5.75), Inches(0.36),
                        "  ▸  " + item, size=11, color=C_LIGHT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14  —  TESTING
# ════════════════════════════════════════════════════════════════════════════
def slide_14():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Testing")
    slide_title(s, "Testing Strategy & Key Test Cases",
                "Unit tests for AI algorithms + system-level end-to-end tests")
    slide_number(s, 14)

    # Stats row
    for i, (num, label, col) in enumerate([
        ("12", "Unit Tests",   C_ACCENT),
        ("8",  "System Tests", C_GOLD),
        ("100%","Auth Coverage", RGBColor(0x2E, 0x7D, 0x32)),
        ("4",   "AI Modules Tested", RGBColor(0x1A, 0x73, 0xA8)),
    ]):
        sx = Inches(0.35) + i * Inches(3.24)
        stat_box(s, sx, Inches(1.82), Inches(3.0), Inches(1.28), num, label, col)

    # Test case table header
    for j, (header, w, x) in enumerate([
        ("Test ID", 0.8, 0.35),
        ("Module",  2.2, 1.15),
        ("Test Description", 5.5, 3.35),
        ("Expected Result", 4.0, 8.85),
    ]):
        add_rect(s, Inches(x), Inches(3.25), Inches(w), Inches(0.38), fill=C_ACCENT)
        add_textbox(s, Inches(x + 0.05), Inches(3.27), Inches(w - 0.08), Inches(0.32),
                    header, size=11, bold=True, color=C_WHITE, font="Calibri")

    tests = [
        ("UT-08", "Cosine Similarity",    "Two identical event documents compared",              "Score = 1.0"),
        ("UT-09", "Cosine Similarity",    "Two completely unrelated documents",                  "Score ≈ 0.0"),
        ("UT-10", "BM25 Search",          "Query term in only one event, not others",            "That event ranked #1"),
        ("UT-12", "Popularity Score",     "High-rated + featured event vs. low-rated",           "Featured event score higher"),
        ("ST-02", "Smart Search",         "Public user searches 'Dashain festival'",             "Festival events surface first"),
        ("ST-03", "Recommendations",      "User views a festival event, checks sidebar",         "Top 5 similar festival events returned"),
        ("ST-04", "Tag Suggestion",       "Admin types 'Mountain Trekking Festival' in form",    "Tags: trekking, mountain, festival shown"),
    ]

    for i, (tid, module, desc, result) in enumerate(tests):
        bg = C_CARD if i % 2 == 0 else C_BG2
        row_y = Inches(3.65) + i * Inches(0.44)
        for (txt, w, x) in [
            (tid,    0.8, 0.35),
            (module, 2.2, 1.15),
            (desc,   5.5, 3.35),
            (result, 4.0, 8.85),
        ]:
            add_rect(s, Inches(x), row_y, Inches(w), Inches(0.42), fill=bg)
            add_textbox(s, Inches(x + 0.05), row_y + Inches(0.05), Inches(w - 0.08), Inches(0.34),
                        txt, size=10, color=C_LIGHT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15  —  GANTT CHART
# ════════════════════════════════════════════════════════════════════════════
def slide_15():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Timeline")
    slide_title(s, "Project Timeline — Gantt Chart",
                "8-week development plan: Weeks 1-12, Midterm at Week 12")
    slide_number(s, 15)

    weeks = ["W1", "W2", "W3", "W4", "W5", "W6", "W7", "W8", "W9", "W10", "W11", "W12"]
    n = len(weeks)
    chart_x = Inches(3.1)
    chart_y = Inches(1.82)
    cell_w  = (SLIDE_W - chart_x - Inches(0.35)) / n
    row_h   = Inches(0.5)

    # Header row
    add_rect(s, chart_x, chart_y, SLIDE_W - chart_x - Inches(0.35), row_h, fill=C_ACCENT)
    for i, w in enumerate(weeks):
        add_textbox(s, chart_x + i * cell_w, chart_y + Inches(0.1),
                    cell_w, Inches(0.3), w, size=10, bold=True,
                    color=C_WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Midterm marker
    marker_x = chart_x + 11 * cell_w + cell_w / 2
    add_rect(s, marker_x - Inches(0.015), chart_y, Inches(0.03),
             Inches(0.5 + 7 * 0.5), fill=C_GOLD)
    add_textbox(s, marker_x - Inches(0.5), chart_y + Inches(0.5 * 7) + Inches(0.05),
                Inches(1.1), Inches(0.3), "MIDTERM", size=8, bold=True,
                color=C_GOLD, align=PP_ALIGN.CENTER, font="Calibri")

    tasks = [
        ("Project Setup & DB Schema",          0, 2,  C_ACCENT),
        ("Auth Module (JWT + Roles)",           1, 3,  RGBColor(0xAD, 0x1D, 0x38)),
        ("Event CRUD API",                      2, 5,  RGBColor(0x6A, 0x0F, 0x1F)),
        ("Admin Dashboard UI",                 3, 6,  C_GOLD),
        ("Calendar Component",                 4, 7,  RGBColor(0xB8, 0x96, 0x20)),
        ("Public Portal & Filters",            5, 8,  RGBColor(0x2E, 0x7D, 0x32)),
        ("AI: TF-IDF + Cosine Similarity",     7, 11, RGBColor(0x1A, 0x73, 0xA8)),
        ("AI: BM25 Search Ranking",            8, 11, RGBColor(0x13, 0x5A, 0x8A)),
        ("AI: Tag Suggestion",                 8, 11, RGBColor(0x0D, 0x47, 0x70)),
        ("AI: Popularity Scoring",             6, 10, RGBColor(0x00, 0x5C, 0x80)),
        ("Testing & Documentation",            9, 12, RGBColor(0x4A, 0x14, 0x5E)),
    ]

    for i, (task, start, end, col) in enumerate(tasks):
        ty = chart_y + (i + 1) * row_h
        bg = C_CARD if i % 2 == 0 else C_BG2
        # row bg
        add_rect(s, Inches(0.35), ty, chart_x - Inches(0.35), row_h, fill=bg)
        add_textbox(s, Inches(0.4), ty + Inches(0.1), chart_x - Inches(0.5), row_h - Inches(0.1),
                    task, size=9.5, color=C_LIGHT, font="Calibri")
        # grid bg
        add_rect(s, chart_x, ty, SLIDE_W - chart_x - Inches(0.35), row_h, fill=bg)
        # bar
        bx = chart_x + start * cell_w + Inches(0.04)
        bw = (end - start) * cell_w - Inches(0.08)
        add_rect(s, bx, ty + Inches(0.1), bw, row_h - Inches(0.2), fill=col)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16  —  CHALLENGES & SOLUTIONS
# ════════════════════════════════════════════════════════════════════════════
def slide_16():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Learnings")
    slide_title(s, "Challenges Faced & How We Solved Them")
    slide_number(s, 16)

    challenges = [
        (
            "Dual Calendar Support (AD + BS)",
            "Bikram Sambat has variable month lengths and no standard library support in .NET",
            "Implemented a custom date conversion utility in both C# and TypeScript that maps AD↔BS using a lookup table for all year ranges covered by the event calendar",
        ),
        (
            "Stale Entity Error in EF Core",
            "Reading with AsNoTracking() then immediately calling Update() caused entity-already-tracked conflicts",
            "Solved by ensuring all GetById calls use AsNoTracking(), then Update() re-attaches the detached entity cleanly — no SaveChanges twice",
        ),
        (
            "AI Cold-Start Problem",
            "Content-based filtering needs a corpus of events to build TF-IDF vectors — an empty database returns nothing useful",
            "Seeded the database with 15+ real NTB events at startup so the recommendation engine has a meaningful corpus from day one",
        ),
        (
            "Dynamic SQL vs SQL Injection",
            "The filter query needed dynamic WHERE clauses but building SQL strings from user input is dangerous",
            "Used Dapper DynamicParameters so user inputs are always parameterised — the query structure is code-controlled, the values are always bound",
        ),
    ]

    for i, (title, challenge, solution) in enumerate(challenges):
        r = i // 2
        c = i % 2
        x = Inches(0.35) + c * Inches(6.5)
        y = Inches(1.85) + r * Inches(2.55)
        add_rect(s, x, y, Inches(6.2), Inches(2.35), fill=C_CARD,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_rect(s, x, y, Inches(6.2), Inches(0.055), fill=C_ACCENT)
        add_textbox(s, x + Inches(0.18), y + Inches(0.12), Inches(5.8), Inches(0.38),
                    title, size=13, bold=True, color=C_WHITE, font="Calibri")
        # Problem label
        add_rect(s, x + Inches(0.18), y + Inches(0.58), Inches(0.88), Inches(0.25),
                 fill=RGBColor(0x6A, 0x1A, 0x1A))
        add_textbox(s, x + Inches(0.18), y + Inches(0.58), Inches(0.9), Inches(0.24),
                    "Problem", size=8, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER, font="Calibri")
        add_textbox(s, x + Inches(0.18), y + Inches(0.86), Inches(5.75), Inches(0.38),
                    challenge, size=10, color=RGBColor(0xFF, 0xAA, 0xAA),
                    wrap=True, font="Calibri")
        # Solution label
        add_rect(s, x + Inches(0.18), y + Inches(1.28), Inches(0.88), Inches(0.25),
                 fill=RGBColor(0x1B, 0x5E, 0x20))
        add_textbox(s, x + Inches(0.18), y + Inches(1.28), Inches(0.9), Inches(0.24),
                    "Solution", size=8, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER, font="Calibri")
        add_textbox(s, x + Inches(0.18), y + Inches(1.56), Inches(5.75), Inches(0.65),
                    solution, size=10, color=RGBColor(0xAA, 0xFF, 0xBB),
                    wrap=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17  —  CONCLUSION + NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
def slide_17():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)
    gold_bar(s, y=0, h=Inches(0.055))
    footer_line(s)
    section_tag(s, "Conclusion")
    slide_title(s, "Conclusion & What's Next",
                "Midterm checkpoint — core platform complete, AI in final integration")
    slide_number(s, 17)

    # Done column
    add_rect(s, Inches(0.35), Inches(1.82), Inches(6.1), Inches(0.44),
             fill=RGBColor(0x1B, 0x5E, 0x20))
    add_textbox(s, Inches(0.55), Inches(1.84), Inches(5.8), Inches(0.38),
                "What We've Achieved", size=14, bold=True, color=C_WHITE, font="Calibri")

    achieved = [
        "Complete event management backend — Clean Architecture, EF Core + Dapper, PostgreSQL",
        "Role-based authentication — JWT with refresh token rotation, BCrypt password hashing",
        "Full admin portal — events, categories, tags, users, lifecycle management",
        "Responsive public portal — homepage, calendar, list view, filtering",
        "Dual-calendar (AD + BS) support throughout the entire application",
        "4 AI module architectures designed and implementation in progress",
    ]
    for i, item in enumerate(achieved):
        y = Inches(2.3) + i * Inches(0.72)
        add_rect(s, Inches(0.35), y, Inches(6.1), Inches(0.62), fill=C_CARD)
        add_textbox(s, Inches(0.55), y + Inches(0.1), Inches(5.7), Inches(0.44),
                    "  ✓  " + item, size=11, color=C_LIGHT, wrap=True, font="Calibri")

    # Next steps column
    add_rect(s, Inches(6.75), Inches(1.82), Inches(6.2), Inches(0.44),
             fill=RGBColor(0x0D, 0x47, 0xA1))
    add_textbox(s, Inches(6.95), Inches(1.84), Inches(5.9), Inches(0.38),
                "Final Defense Targets", size=14, bold=True, color=C_WHITE, font="Calibri")

    next_steps = [
        ("Complete AI Integration",   "Wire all 4 AI modules to live API endpoints and frontend UI"),
        ("Recommendation API",        "GET /api/events/{id}/recommendations live and tested"),
        ("BM25 Search Live",          "Replace current ILIKE search with BM25-ranked results"),
        ("Tag Suggestion UI",         "Chip suggestions visible in admin event form in real time"),
        ("Comprehensive Testing",     "Full unit + system test suite, result analysis report"),
        ("Deployment & Demo",         "Dockerized deployment, live demo environment ready"),
    ]
    for i, (title, detail) in enumerate(next_steps):
        y = Inches(2.3) + i * Inches(0.72)
        add_rect(s, Inches(6.75), y, Inches(6.2), Inches(0.62), fill=C_CARD,
                 line_color=C_DIVIDER, line_w=Pt(1))
        add_rect(s, Inches(6.75), y, Inches(0.055), Inches(0.62), fill=C_GOLD)
        add_textbox(s, Inches(6.95), y + Inches(0.04), Inches(5.8), Inches(0.26),
                    title, size=12, bold=True, color=C_GOLD, font="Calibri")
        add_textbox(s, Inches(6.95), y + Inches(0.32), Inches(5.8), Inches(0.26),
                    detail, size=10, color=C_MUTED, font="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18  —  THANK YOU / Q&A
# ════════════════════════════════════════════════════════════════════════════
def slide_18():
    s = prs.slides.add_slide(blank_layout)
    full_bg(s)

    # Left strip
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, fill=C_GOLD)

    # Decorative shapes
    shape = s.shapes.add_shape(9, SLIDE_W - Inches(4), -Inches(1.5), Inches(5.5), Inches(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x18, 0x40, 0x28)
    shape.line.fill.background()

    shape2 = s.shapes.add_shape(9, Inches(0.5), SLIDE_H - Inches(2.8), Inches(3), Inches(3))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0x15, 0x33, 0x20)
    shape2.line.fill.background()

    add_rect(s, 0, 0, SLIDE_W, Inches(0.055), fill=C_GOLD)

    # Main
    add_textbox(s, Inches(0.6), Inches(1.6), Inches(9), Inches(1.1),
                "Thank You", size=56, bold=True, color=C_WHITE, font="Calibri")

    add_rect(s, Inches(0.6), Inches(2.82), Inches(2.8), Inches(0.055), fill=C_ACCENT)

    add_textbox(s, Inches(0.6), Inches(3.0), Inches(8.5), Inches(0.5),
                "Questions & Discussion Welcome", size=20, italic=True,
                color=C_MUTED, font="Calibri")

    # Info row
    for i, (label, value, col) in enumerate([
        ("Project",   "NTB Event Management System",          C_WHITE),
        ("Developer", "Anjal Joshi  •  anjaljoshi6@gmail.com", C_MUTED),
        ("Defense",   "Midterm  •  June 2026  •  Semester VIII", C_MUTED),
    ]):
        y = Inches(3.82) + i * Inches(0.52)
        add_textbox(s, Inches(0.6), y, Inches(1.6), Inches(0.44),
                    label + ":", size=12, bold=True, color=C_GOLD, font="Calibri")
        add_textbox(s, Inches(2.3), y, Inches(7.5), Inches(0.44),
                    value, size=12, color=col, font="Calibri")

    # Tech stack chips bottom
    for i, (tag, col) in enumerate([
        ("ASP.NET Core 8", C_ACCENT),
        ("SvelteKit",       C_GOLD),
        ("PostgreSQL",      RGBColor(0x33, 0x6F, 0x99)),
        ("TF-IDF + BM25",   RGBColor(0x2E, 0x7D, 0x32)),
        ("Clean Architecture", RGBColor(0x7B, 0x1F, 0xA2)),
    ]):
        bx = Inches(0.6) + i * Inches(2.45)
        add_rect(s, bx, Inches(5.6), Inches(2.3), Inches(0.42), fill=C_CARD,
                 line_color=col, line_w=Pt(1.5))
        add_textbox(s, bx + Inches(0.1), Inches(5.62), Inches(2.1), Inches(0.36),
                    tag, size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Bottom bar
    add_rect(s, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), fill=C_CARD)
    add_textbox(s, Inches(0.6), SLIDE_H - Inches(0.35), Inches(8), Inches(0.32),
                "Department of Computer Science  •  Pokhara University  •  CACS452 Project III",
                size=10, color=C_MUTED)
    slide_number(s, 18)


# ════════════════════════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ════════════════════════════════════════════════════════════════════════════
slide_01()   # Cover
slide_02()   # Agenda
slide_03()   # Problem Statement
slide_04()   # Objectives
slide_05()   # Architecture
slide_06()   # Tech Stack
slide_07()   # Progress
slide_08()   # Data Model
slide_09()   # AI Overview
slide_10()   # AI: TF-IDF + Cosine
slide_11()   # AI: BM25
slide_12()   # AI: Tag Suggestion + Popularity
slide_13()   # Demo Screens
slide_14()   # Testing
slide_15()   # Gantt Chart
slide_16()   # Challenges
slide_17()   # Conclusion
slide_18()   # Thank You

# ════════════════════════════════════════════════════════════════════════════
out = "/home/notcool/Desktop/ntb-event/NTB_Event_MidTerm_Defense_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
